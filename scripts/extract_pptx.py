"""
需求文档提取脚本 — 将每页幻灯片导出为完整 PNG 截图

用法：
    python extract_pptx.py --input requirements.pptx [--slides 1-3,5,7-9] [--outdir pptx_output]
    python extract_pptx.py --input requirements.pdf  [--slides 1-3,5,7-9] [--outdir pptx_output]

输出：
    outdir/ 下每页一张 PNG 图片（如 slide_1.png, slide_2.png），供 LLM 用 Read 工具直接查看

依赖：
    - pip: PyMuPDF (fitz)
    - 仅 PPTX 输入时需要系统安装 LibreOffice（用于 pptx → pdf 转换）
"""

import argparse
import os
import re
import subprocess
import sys


def parse_slide_range(spec, total):
    if not spec:
        return set(range(1, total + 1))
    result = set()
    for part in spec.split(','):
        part = part.strip()
        m = re.match(r'^(\d+)-(\d+)$', part)
        if m:
            lo, hi = int(m.group(1)), int(m.group(2))
            result.update(range(lo, min(hi, total) + 1))
        elif part.isdigit():
            n = int(part)
            if 1 <= n <= total:
                result.add(n)
    return result


def find_libreoffice():
    # Windows 上 soffice.exe 是 GUI 启动器，subprocess 调用 --version 会弹
    # "Press Enter to continue" 控制台窗口并阻塞。soffice.com 是控制台变体，
    # 静默运行，必须优先匹配 .com。
    if sys.platform == 'win32':
        candidates = []
        for base in [os.environ.get('PROGRAMFILES', r'C:\Program Files'),
                     os.environ.get('PROGRAMFILES(X86)', r'C:\Program Files (x86)')]:
            candidates.append(os.path.join(base, 'LibreOffice', 'program', 'soffice.com'))
        candidates += ['soffice.com', 'soffice', 'libreoffice']
    else:
        candidates = ['libreoffice', 'soffice']
    for cmd in candidates:
        try:
            subprocess.run([cmd, '--version'], capture_output=True, timeout=10)
            return cmd
        except (FileNotFoundError, subprocess.TimeoutExpired, OSError):
            continue
    return None


def pptx_to_pdf(pptx_path, outdir):
    lo = find_libreoffice()
    if not lo:
        print("ERROR: 未找到 LibreOffice。请安装 LibreOffice：https://www.libreoffice.org/download/")
        sys.exit(1)

    print(f"📦 正在转换 PPTX → PDF（LibreOffice）...")
    result = subprocess.run(
        [lo, '--headless', '--convert-to', 'pdf', '--outdir', outdir, pptx_path],
        capture_output=True, text=True, timeout=120
    )
    if result.returncode != 0:
        print(f"ERROR: LibreOffice 转换失败：{result.stderr}")
        sys.exit(1)

    pdf_name = os.path.splitext(os.path.basename(pptx_path))[0] + '.pdf'
    pdf_path = os.path.join(outdir, pdf_name)
    if not os.path.exists(pdf_path):
        print(f"ERROR: PDF 文件未生成：{pdf_path}")
        sys.exit(1)

    return pdf_path


def pdf_to_images(pdf_path, outdir, slides_spec, dpi, cleanup_pdf=False):
    import fitz
    doc = fitz.open(pdf_path)
    total = len(doc)
    selected = parse_slide_range(slides_spec, total)

    exported = []
    for i in range(total):
        page_num = i + 1
        if page_num not in selected:
            continue
        page = doc[i]
        pix = page.get_pixmap(dpi=dpi)
        out_path = os.path.join(outdir, f'slide_{page_num}.png')
        pix.save(out_path)
        exported.append(out_path)
        print(f"📄 Slide {page_num}/{total} → {out_path}")

    doc.close()
    if cleanup_pdf:
        os.remove(pdf_path)

    return exported, total


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument('--input', required=True, help='PPTX 或 PDF 文件路径')
    parser.add_argument('--slides', default=None, help='页码范围，如 1-3,5,7-9')
    parser.add_argument('--outdir', default='pptx_output', help='图片输出目录')
    parser.add_argument('--dpi', type=int, default=200, help='导出图片 DPI（默认 200）')
    parser.add_argument('--info', action='store_true', help='仅输出页数信息，不执行导出')
    args = parser.parse_args()

    input_path = os.path.abspath(args.input)
    if not os.path.exists(input_path):
        print(f"ERROR: 文件不存在：{input_path}")
        sys.exit(1)

    ext = os.path.splitext(input_path)[1].lower()

    # --info 模式
    if args.info:
        if ext == '.pdf':
            import fitz
            doc = fitz.open(input_path)
            total = len(doc)
            doc.close()
        else:
            from pptx import Presentation
            prs = Presentation(input_path)
            total = len(prs.slides)
        print(f"📊 共 {total} 页")
        sys.exit(0)

    os.makedirs(args.outdir, exist_ok=True)
    outdir = os.path.abspath(args.outdir)

    if ext == '.pdf':
        # PDF 直接渲染，无需 LibreOffice
        print(f"📄 输入为 PDF，直接渲染图片...")
        exported, total = pdf_to_images(input_path, outdir, args.slides, args.dpi, cleanup_pdf=False)
    else:
        # PPTX 先转 PDF，再渲染
        pdf_path = pptx_to_pdf(input_path, outdir)
        exported, total = pdf_to_images(pdf_path, outdir, args.slides, args.dpi, cleanup_pdf=True)

    print(f"\n✅ 导出完成：{len(exported)}/{total} 页")
    print(f"📁 图片目录：{outdir}")
    for p in exported:
        print(f"  - {p}")


if __name__ == '__main__':
    main()
